import openai
from newspaper import Article
from pptx import Presentation
from pptx.util import Inches

# STEP 1: Insert your OpenAI API key
openai.api_key = "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"

# STEP 2: Function to scrape article from URL
def scrape_article(url):
    from newspaper import Article
    article = Article(url)
    try:
        article.download()
        article.parse()
        return article.title, article.text
    except Exception as e:
        print(f"Scraping failed: {e}")
        print("Paste the article content manually below ⬇️:")
        text = input("Paste article text: ")
        title = input("Enter a title for the presentation: ")
        return title, text

# STEP 3: Ask GPT to turn text into slide content
from openai import OpenAI

client = OpenAI()

def generate_slide_content(text):
    prompt = f"""Break down the following article into a PowerPoint outline with 5–10 slides.
Each slide should have a title and 3–5 bullet points summarizing the content.

Article:
{text}

Format:
Slide 1 Title: ...
- Bullet 1
- Bullet 2
Slide 2 Title: ...
"""
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=1000
    )

    return response.choices[0].message.content



# STEP 4: Generate PowerPoint file
def create_ppt_from_gpt(gpt_response, output_filename="presentation.pptx"):
    prs = Presentation()
    slides_data = gpt_response.strip().split("Slide ")[1:]

    for slide_raw in slides_data:
        lines = slide_raw.strip().splitlines()
        if not lines:
            continue
        title_line = lines[0].replace("Title:", "").strip()
        content_lines = [line.strip("- ").strip() for line in lines[1:] if line.startswith("-")]

        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_line
        content = slide.placeholders[1]
        content.text = "\n".join(content_lines)

    prs.save(output_filename)
    return output_filename

# STEP 5: Run it on an example article
if __name__ == "__main__":
    url = input("Paste the article URL: ")
    print("Scraping article...")
    title, text = scrape_article(url)
    print("Generating slide content with GPT...")
    gpt_output = generate_slide_content(text)
    print("Creating PowerPoint...")
    filename = create_ppt_from_gpt(gpt_output)
    print(f"Done! Your file is saved as: {filename}")



    